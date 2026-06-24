"""Local document extraction for the no-API MVP.

This module intentionally performs no network calls. It extracts text from the
trainer's uploaded file so the prototype analysis engine can operate entirely
on the local machine.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md"}


class UnsupportedDocumentError(ValueError):
    """Raised when the uploaded file is outside the MVP's supported formats."""


def extract_document_text(filename: str, content: bytes) -> dict:
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise UnsupportedDocumentError(
            f"Unsupported file type '{extension}'. Use PDF, PPTX, DOCX, TXT or MD."
        )

    if extension == ".pdf":
        return _extract_pdf(content)
    if extension == ".pptx":
        return _extract_pptx(content)
    if extension == ".docx":
        return _extract_docx(content)
    return _extract_text(content)


def _extract_pdf(content: bytes) -> dict:
    reader = PdfReader(BytesIO(content))
    units = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            units.append({"label": f"Page {index}", "text": text})
    return {"kind": "PDF", "unit_count": len(reader.pages), "units": units}


def _extract_pptx(content: bytes) -> dict:
    presentation = Presentation(BytesIO(content))
    units = []
    for index, slide in enumerate(presentation.slides, start=1):
        fragments = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                fragments.append(text.strip())
        if fragments:
            units.append({"label": f"Slide {index}", "text": "\n".join(fragments)})
    return {"kind": "PPTX", "unit_count": len(presentation.slides), "units": units}


def _extract_docx(content: bytes) -> dict:
    document = Document(BytesIO(content))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]

    # Table text often contains the most useful procedure and comparison data.
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    return {
        "kind": "DOCX",
        "unit_count": max(1, len(document.sections)),
        "units": [{"label": "Document", "text": "\n".join(paragraphs)}],
    }


def _extract_text(content: bytes) -> dict:
    text = content.decode("utf-8", errors="replace")
    return {"kind": "Text", "unit_count": 1, "units": [{"label": "Document", "text": text}]}
